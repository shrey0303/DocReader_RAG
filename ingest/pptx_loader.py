import logging
from typing import List
from pptx import Presentation

logging.basicConfig(level=logging.INFO)

def load_and_chunk_pptx(path: str, chunk_size: int = 500) -> List[str]:
    """
    Reads a PPTX file and splits its text into chunks.
    Args:
        path: Path to the PPTX file.
        chunk_size: Number of characters per chunk.
    Returns:
        List of text chunks.
    """
    prs = Presentation(path)
    text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    logging.info(f"Loaded {len(chunks)} chunks from {path}")
    return chunks 